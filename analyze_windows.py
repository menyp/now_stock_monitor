import os
import pandas as pd
import yfinance as yf
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from matplotlib.patches import Rectangle
from datetime import datetime, timedelta
import numpy as np
import seaborn as sns
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from collections import defaultdict, Counter
import traceback

def run_analysis_for_years(YEARS=10):
    print(f"Running analysis for {YEARS} years...")
    # --- Configuration ---
    TICKER = "NOW"
    WINDOW_MONTHS = [2, 5, 8, 11]  # Feb, May, Aug, Nov (employee trading windows)
    WINDOW_LABELS = ["Q1-Feb", "Q2-May", "Q3-Aug", "Q4-Nov"]
    PLOTS_DIR = os.path.join(os.path.dirname(__file__), "plots")
    os.makedirs(PLOTS_DIR, exist_ok=True)
    LOGO_PATH = os.path.join(os.path.dirname(__file__), "logo.png")

    # --- Download Data ---
    end_date = datetime.today()
    start_date = end_date - timedelta(days=YEARS * 365)
    df = yf.download(TICKER, start=start_date.strftime("%Y-%m-%d"), end=end_date.strftime("%Y-%m-%d"))
    df.reset_index(inplace=True)

    # --- Prepare Results ---
    results = []
    window_plot_files = []
    # --- For optimal selling window analysis ---
    best_windows = []
    window_lengths = [3, 4, 5]
    for year in range(end_date.year - YEARS + 1, end_date.year + 1):
        for i, month in enumerate(WINDOW_MONTHS):
            window_mask = (df['Date'] >= pd.Timestamp(year=year, month=month, day=1)) & \
                          (df['Date'] < pd.Timestamp(year=year if month < 12 else year+1, month=month % 12 + 1, day=1))
            window_df = df.loc[window_mask].copy()
            if window_df.empty or len(window_df) < 3:
                continue

            window_df = window_df.reset_index(drop=True)
            first_day = window_df.iloc[0]
            last_day = window_df.iloc[-1]
            max_day = window_df.loc[window_df['Close'].idxmax()]
            min_day = window_df.loc[window_df['Close'].idxmin()]

            total_days = len(window_df)
            max_idx = window_df['Close'].idxmax()
            if isinstance(max_idx, (pd.Series, list, tuple, np.ndarray)):
                max_pos = int(max_idx.iloc[0] if hasattr(max_idx, 'iloc') else max_idx[0])
            else:
                max_pos = int(max_idx)

            if max_pos == 0:
                max_when = "Start"
            elif max_pos == total_days - 1:
                max_when = "End"
            else:
                max_when = "Middle"

            def get_date(val):
                dt = pd.to_datetime(val)
                if isinstance(dt, pd.Series):
                    dt = dt.iloc[0]
                return dt.strftime("%Y-%m-%d")

            def get_scalar(val):
                if isinstance(val, pd.DataFrame):
                    return float(val.iloc[0, 0])
                if isinstance(val, pd.Series):
                    return float(val.iloc[0])
                return float(val)

            results.append({
                "Window": f"{year}-{WINDOW_LABELS[i]}",
                "First_Date": get_date(first_day['Date']),
                "First_Close": get_scalar(first_day['Close']),
                "Last_Date": get_date(last_day['Date']),
                "Last_Close": get_scalar(last_day['Close']),
                "Max_Date": get_date(max_day['Date']),
                "Max_Close": get_scalar(max_day['Close']),
                "Max_When": max_when,
                "Min_Date": get_date(min_day['Date']),
                "Min_Close": get_scalar(min_day['Close']),
            })

            # Save improved individual window plot (not returned to UI for now)
            sns.set_theme(style="whitegrid")
            fig, ax = plt.subplots(figsize=(6, 3))
            shadow = Rectangle((-0.03, -0.03), 1.06, 1.06, transform=ax.transAxes,
                               color='grey', alpha=0.13, zorder=1, linewidth=0)
            ax.add_patch(shadow)
            ax.set_facecolor('#f5f9ff')
            fig.patch.set_facecolor('#f5f9ff')
            ax.plot(window_df['Date'], window_df['Close'], marker='o', color='tab:blue', linewidth=2, markersize=6, label='Close Price', zorder=3)
            ax.scatter([max_day['Date']], [max_day['Close']], color='tab:green', s=80, edgecolor='black', zorder=4, label='Max')
            ax.scatter([min_day['Date']], [min_day['Close']], color='tab:red', s=80, edgecolor='black', zorder=4, label='Min')
            ax.scatter([first_day['Date']], [first_day['Close']], color='tab:cyan', s=80, edgecolor='black', zorder=4, label='First')
            ax.scatter([last_day['Date']], [last_day['Close']], color='tab:orange', s=80, edgecolor='black', zorder=4, label='Last')
            offsets = {
                'Max': (0, -35),
                'Min': (0, 30),
                'First': (-55, 0),
                'Last': (55, 0)
            }
            for pt, color, txt in zip(
                [max_day, min_day, first_day, last_day],
                ['tab:green', 'tab:red', 'tab:cyan', 'tab:orange'],
                ['Max', 'Min', 'First', 'Last']
            ):
                close_val = extract_scalar(pt['Close'])
                date_val = extract_date(pt['Date'])
                dx, dy = offsets[txt]
                ax.annotate(
                    f"{txt}\n{close_val:.2f}",
                    (date_val, close_val),
                    textcoords="offset points",
                    xytext=(dx, dy),
                    ha='center',
                    fontsize=9,
                    color=color,
                    fontweight='bold',
                    bbox=dict(boxstyle="round,pad=0.2", fc="white", ec=color, lw=1, alpha=0.7),
                    arrowprops=dict(arrowstyle="->", lw=1, color=color, shrinkA=0, shrinkB=5),
                    zorder=10
                )
            ax.grid(True, linestyle='-', linewidth=0.5, color='#b6c6e3', alpha=0.6, zorder=2)
            ax.xaxis.set_major_locator(mdates.AutoDateLocator())
            ax.xaxis.set_major_formatter(mdates.DateFormatter('%b %d'))
            fig.autofmt_xdate(rotation=45)
            ax.set_xlabel("Date", fontsize=12)
            ax.set_ylabel("Close Price ($)", fontsize=12)
            ax.set_title(f"{TICKER} {year} {WINDOW_LABELS[i]} Window", fontsize=16, fontweight='bold', pad=16)
            fig.tight_layout()
            plot_file = os.path.join(PLOTS_DIR, f"window_plot_{year}_{WINDOW_LABELS[i]}.png")
            plt.savefig(plot_file, dpi=150)
            plt.close()
            window_plot_files.append((f"{year}-{WINDOW_LABELS[i]}", os.path.abspath(plot_file)))

            n_days = len(window_df)
            for win_len in window_lengths:
                if n_days < win_len:
                    continue
                best_avg = -np.inf
                best_start = None
                for start in range(n_days - win_len + 1):
                    avg = float(window_df['Close'].iloc[start:start+win_len].mean())
                    if avg > best_avg:
                        best_avg = avg
                        best_start = start
                best_windows.append((win_len, best_start, best_avg))

    # --- Output Results ---
    summary_df = pd.DataFrame(results)
    max_when_counts = summary_df['Max_When'].value_counts()
    total_windows = max_when_counts.sum()
    periods = ["Start", "Middle", "End"]
    most_common_period = max_when_counts.idxmax()
    if most_common_period == "Start":
        reco = "Based on the last {} years, the highest price during trading windows most often occurs at the START of the window. Consider selling early in the window.".format(YEARS)
    elif most_common_period == "End":
        reco = "Based on the last {} years, the highest price during trading windows most often occurs at the END of the window. Consider selling late in the window.".format(YEARS)
    else:
        reco = "Based on the last {} years, the highest price during trading windows most often occurs in the MIDDLE of the window. Consider monitoring prices throughout the window and selling when a spike occurs.".format(YEARS)

    # --- Composite Visualization ---
    summary_df['Window_Year'] = summary_df['Window'].str[:4]
    summary_df['Quarter'] = summary_df['Window'].str[-3:]
    period_map = {'Start': 0, 'Middle': 1, 'End': 2}
    summary_df['Max_When_Num'] = summary_df['Max_When'].map(period_map)

    fig, axs = plt.subplots(2, 2, figsize=(12, 7), gridspec_kw={'height_ratios': [2, 1]})
    counts = [max_when_counts.get(period, 0) for period in periods]
    colors = ['tab:green' if period == most_common_period else 'tab:gray' for period in periods]
    bars = axs[0, 0].bar(periods, counts, color=colors, edgecolor='black')
    axs[0, 0].set_title("When Does the Highest Price Occur?")
    axs[0, 0].set_ylabel("Number of Windows")
    for bar, count in zip(bars, counts):
        axs[0, 0].text(bar.get_x() + bar.get_width()/2, bar.get_height(), str(count), ha='center', va='bottom', fontsize=12)
    highlight_patch = mpatches.Patch(color='tab:green', label='Recommended Period')
    axs[0, 0].legend(handles=[highlight_patch])
    sns.stripplot(
        x='Max_When', y='Window_Year', data=summary_df,
        order=periods, ax=axs[0, 1], size=10, palette='Set2', jitter=True
    )
    axs[0, 1].set_title("Max Price Timing by Window")
    axs[0, 1].set_xlabel("Period in Window")
    axs[0, 1].set_ylabel("Year")
    axs[1, 0].axis('off')
    axs[1, 0].text(0, 0.5, reco, fontsize=13, bbox=dict(facecolor='lightyellow', edgecolor='black'))
    axs[1, 1].axis('off')
    table_data = [[period, counts[i]] for i, period in enumerate(periods)]
    table = axs[1, 1].table(cellText=table_data, colLabels=["Period", "Count"], loc='center')
    table.auto_set_font_size(False)
    table.set_fontsize(12)
    table.scale(1.2, 1.2)
    plt.tight_layout()
    composite_path = os.path.join(PLOTS_DIR, "composite_summary.png")
    plt.savefig(composite_path, dpi=300)
    # Convert composite plot to base64
    import base64
    from io import BytesIO
    buf = BytesIO()
    plt.savefig(buf, format='png')
    buf.seek(0)
    composite_base64 = base64.b64encode(buf.read()).decode('utf-8')
    plt.close()

    # --- Bar Chart Per Window (Quarter) ---
    window_charts = {}
    for i, label in enumerate(WINDOW_LABELS):
        window_mask = summary_df['Window'].str.contains(label)
        window_df = summary_df[window_mask]
        counts = window_df['Max_When'].value_counts()
        periods = ["Start", "Middle", "End"]
        values = [counts.get(p, 0) for p in periods]
        plt.figure(figsize=(5, 4))
        bars = plt.bar(periods, values, color=['tab:blue', 'tab:orange', 'tab:green'])
        plt.title(f"Max Price Timing: {label}")
        plt.ylabel("Number of Years")
        for bar, val in zip(bars, values):
            plt.text(bar.get_x() + bar.get_width()/2, bar.get_height(), str(val), ha='center', va='bottom', fontsize=12)
        plt.tight_layout()
        # Save and encode as base64
        buf = BytesIO()
        plt.savefig(buf, format='png')
        buf.seek(0)
        window_charts[label] = base64.b64encode(buf.read()).decode('utf-8')
        plt.close()

    # --- Gather per-year, per-window daily close data ---
    per_year_window_data = {}
    for year in range(end_date.year - YEARS + 1, end_date.year + 1):
        year_str = str(year)
        per_year_window_data[year_str] = {}
        for i, month in enumerate(WINDOW_MONTHS):
            label = WINDOW_LABELS[i]
            window_mask = (df['Date'] >= pd.Timestamp(year=year, month=month, day=1)) & \
                          (df['Date'] < pd.Timestamp(year=year if month < 12 else year+1, month=month % 12 + 1, day=1))
            window_df = df.loc[window_mask].copy()
            if window_df.empty:
                continue
            per_year_window_data[year_str][label] = [
                {'date': d.strftime('%Y-%m-%d'), 'close': float(c)}
                for d, c in zip(window_df['Date'], window_df['Close'])
            ]

    # --- Aggregate all years for each window ---
    per_window_all_years = {label: {} for label in WINDOW_LABELS}
    for year in range(end_date.year - YEARS + 1, end_date.year + 1):
        year_str = str(year)
        for i, month in enumerate(WINDOW_MONTHS):
            label = WINDOW_LABELS[i]
            window_mask = (df['Date'] >= pd.Timestamp(year=year, month=month, day=1)) & \
                          (df['Date'] < pd.Timestamp(year=year if month < 12 else year+1, month=month % 12 + 1, day=1))
            window_df = df.loc[window_mask].copy()
            if window_df.empty:
                continue
            per_window_all_years[label][year_str] = [
                {'date': d.strftime('%Y-%m-%d'), 'close': float(c)}
                for d, c in zip(window_df['Date'], window_df['Close'])
            ]

    return {
        'recommendation': reco,
        'summary': summary_df.to_dict(orient='records'),
        'composite_chart': composite_base64,
        'window_charts': window_charts,
        'per_year_window_data': per_year_window_data,
        'per_window_all_years': per_window_all_years,
    }

if __name__ == "__main__":
    # Default: run for 10 years and print results
    result = run_analysis_for_years(10)
    print("Analysis complete. See output files and plots directory.")

def extract_scalar(val):
    if isinstance(val, pd.DataFrame):
        return float(val.iloc[0, 0])
    elif isinstance(val, pd.Series):
        return float(val.iloc[0])
    else:
        return float(val)

def extract_date(val):
    if isinstance(val, pd.DataFrame):
        return val.iloc[0, 0]
    elif isinstance(val, pd.Series):
        return val.iloc[0]
    else:
        return val

# (All analysis logic is now inside run_analysis_for_years)

def remove_footers(slide):
    for shape in list(slide.shapes):
        if not shape.is_placeholder:
            continue
        phf = shape.placeholder_format
        if phf.type in (PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.SLIDE_NUMBER):
            sp = shape
            slide.shapes._spTree.remove(sp._element)

def add_logo(slide):
    logo_width = Inches(1.0)
    logo_height = Inches(0.5)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    left = slide_width - logo_width - Inches(0.2)
    top = slide_height - logo_height - Inches(0.2)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, left, top, width=logo_width, height=logo_height)
    else:
        logo_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, logo_width, logo_height)
        logo_shape.fill.solid()
        logo_shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
        logo_shape.line.color.rgb = RGBColor(120, 120, 120)
        tf = logo_shape.text_frame
        tf.text = "LOGO"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

print("\n--- Starting PowerPoint Export ---")
try:
    print(f"Number of results: {len(results)}")
    print(f"Number of window_plot_files: {len(window_plot_files)}")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 249, 255)
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(9), Inches(1.5))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(0, 70, 122)
    title_shape.line.fill.background()
    title_tf = title_shape.text_frame
    title_tf.text = "ServiceNow (NOW) Employee Trading Window Analysis"
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_tf.paragraphs[0].font.alignment = PP_ALIGN.CENTER
    subtitle_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.7), Inches(9), Inches(0.7))
    subtitle_shape.fill.solid()
    subtitle_shape.fill.fore_color.rgb = RGBColor(255, 255, 204)
    subtitle_shape.line.fill.background()
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.text = f"Analysis period: last {YEARS} years | Generated: {datetime.today().strftime('%Y-%m-%d')}"
    subtitle_tf.paragraphs[0].font.size = Pt(22)
    subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(0, 70, 122)
    subtitle_tf.paragraphs[0].font.alignment = PP_ALIGN.CENTER
    remove_footers(slide)
    add_logo(slide)
    print("Added stylish title slide.")

    img_slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = img_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 249, 255)
    img_title = img_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.8))
    img_title.fill.solid()
    img_title.fill.fore_color.rgb = RGBColor(0, 70, 122)
    img_title.line.fill.background()
    img_title_tf = img_title.text_frame
    img_title_tf.text = "Composite Analysis"
    img_title_tf.paragraphs[0].font.size = Pt(32)
    img_title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    img_title_tf.paragraphs[0].font.alignment = PP_ALIGN.CENTER
    if os.path.exists(os.path.abspath(composite_path)):
        img_slide.shapes.add_picture(os.path.abspath(composite_path), Inches(0.7), Inches(2.0), width=Inches(8.7))
        print(f"Added composite image: {composite_path}")
    else:
        print(f"WARNING: Composite image not found: {composite_path}")
    remove_footers(img_slide)
    add_logo(img_slide)

    reco_slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = reco_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 249, 255)
    card = reco_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(8.5), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 204)
    card.shadow.inherit = False
    card.shadow.blur_radius = 10
    card.shadow.distance = 8
    card.line.color.rgb = RGBColor(0, 70, 122)
    reco_title = reco_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.2), Inches(8.1), Inches(0.8))
    reco_title.fill.solid()
    reco_title.fill.fore_color.rgb = RGBColor(0, 70, 122)
    reco_title.line.fill.background()
    reco_title_tf = reco_title.text_frame
    reco_title_tf.text = "Recommendation"
    reco_title_tf.paragraphs[0].font.size = Pt(32)
    reco_title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    reco_title_tf.paragraphs[0].font.alignment = PP_ALIGN.CENTER
    reco_shape = reco_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.1), Inches(8.1), Inches(2.5))
    reco_shape.fill.solid()
    reco_shape.fill.fore_color.rgb = RGBColor(255, 255, 204)
    reco_shape.line.fill.background()
    reco_tf = reco_shape.text_frame
    reco_tf.text = reco
    reco_tf.paragraphs[0].font.size = Pt(20)
    reco_tf.paragraphs[0].font.color.rgb = RGBColor(0, 70, 122)
    reco_tf.paragraphs[0].font.alignment = PP_ALIGN.LEFT
    remove_footers(reco_slide)
    add_logo(reco_slide)
    print("Added stylish recommendation slide.")

    # --- WHY ARE HIGHS IN THE MIDDLE? INSIGHT SLIDE ---
    why_slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = why_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 249, 255)
    title_shape = why_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(0.7), Inches(7), Inches(1.1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(0, 70, 122)
    title_shape.line.fill.background()
    title_shape.shadow.inherit = False
    title_shape.shadow.blur_radius = 10
    title_shape.shadow.distance = 8
    title_tf = title_shape.text_frame
    title_tf.text = "Why Are Highs in the Middle?"
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_tf.paragraphs[0].font.alignment = PP_ALIGN.CENTER
    card = why_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.1), Inches(7.6), Inches(4.0)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 204)
    card.line.color.rgb = RGBColor(0, 70, 122)
    card.shadow.inherit = False
    card.shadow.blur_radius = 10
    card.shadow.distance = 8
    why_bullets = [
        "• Early in the window, many employees rush to sell, creating initial selling pressure and lower prices.",
        "• As the window progresses, selling pressure subsides and prices often recover or spike.",
        "• New information or market events may emerge mid-window, driving price increases.",
        "• Behavioral factors: employees may 'wait and see,' reducing pressure and allowing a mid-window high.",
        "• Some effect may be due to randomness or mean reversion after early volatility."
    ]
    tf = card.text_frame
    for i, bullet in enumerate(why_bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 70, 122)
        p.level = 0
    remove_footers(why_slide)
    add_logo(why_slide)

    year_to_plots = defaultdict(list)
    for window_title, plot_file in window_plot_files:
        year = window_title.split('-')[0]
        year_to_plots[year].append((window_title, plot_file))

    for year, plots in year_to_plots.items():
        for i in range(0, len(plots), 4):
            group = plots[i:i+4]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 249, 255)
            slide_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            title_tf = slide_title.text_frame
            title_tf.text = f"Employee Trading Window Plots: {year}"
            title_tf.paragraphs[0].font.size = Pt(26)
            title_tf.paragraphs[0].font.bold = True
            title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 70, 122)
            title_tf.paragraphs[0].font.alignment = PP_ALIGN.LEFT
            quarters = ", ".join([plots[j][0].split('-')[1] for j in range(i, min(i+4, len(plots)))])
            slide_subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
            subtitle_tf = slide_subtitle.text_frame
            subtitle_tf.text = f"Quarters: {quarters}"
            subtitle_tf.paragraphs[0].font.size = Pt(16)
            subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(0, 70, 122)
            subtitle_tf.paragraphs[0].font.alignment = PP_ALIGN.LEFT
            positions = [
                (Inches(0.5), Inches(1.5)),
                (Inches(5.0), Inches(1.5)),
                (Inches(0.5), Inches(4.0)),
                (Inches(5.0), Inches(4.0)),
            ]
            for (win_title, plot_file), (left, top) in zip(group, positions):
                abs_plot_file = os.path.abspath(plot_file)
                if not os.path.exists(abs_plot_file):
                    print(f"WARNING: Plot file {abs_plot_file} not found, skipping image.")
                    continue
                slide.shapes.add_picture(abs_plot_file, left, top, width=Inches(PLOT_WIDTH), height=Inches(PLOT_HEIGHT))
            remove_footers(slide)
            add_logo(slide)
            print(f"Added stylish slide for year {year}, plots {i+1} to {i+len(group)}.")

    # --- Key Takeaways Slide (at the end, with improved design) ---
    takeaways_slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = takeaways_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 249, 255)
    title_shape = takeaways_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(0.7), Inches(7), Inches(1.1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(0, 70, 122)
    title_shape.line.fill.background()
    title_shape.shadow.inherit = False
    title_shape.shadow.blur_radius = 10
    title_shape.shadow.distance = 8
    title_tf = title_shape.text_frame
    title_tf.text = "Key Takeaways"
    title_tf.paragraphs[0].font.size = Pt(38)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_tf.paragraphs[0].font.alignment = PP_ALIGN.CENTER
    card = takeaways_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.1), Inches(7.6), Inches(3.6)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 204)
    card.line.color.rgb = RGBColor(0, 70, 122)
    card.shadow.inherit = False
    card.shadow.blur_radius = 10
    card.shadow.distance = 8
    takeaways = [
        f"• Analysis covers {YEARS} years of ServiceNow (NOW) employee trading windows.",
        f"• Highest window price most often occurred at the {most_common_period.upper()} of the window.",
        f"• Number of windows analyzed: {total_windows}",
        f"• Recommendation: {reco}",
        f"• Most effective {best_win_len}-day selling window: days {best_start+1}-{best_start+best_win_len} of each window.",
        "• For best results, monitor trading windows closely and consider timing based on historical trends."
    ]
    tf = card.text_frame
    for i, takeaway in enumerate(takeaways):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = takeaway
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 70, 122)
        p.level = 0
    remove_footers(takeaways_slide)
    add_logo(takeaways_slide)

    pptx_filename = os.path.join(PLOTS_DIR, "ServiceNow_Trading_Window_Analysis.pptx")
    prs.save(pptx_filename)
    print(f"\nPowerPoint presentation saved as {pptx_filename}")

except Exception as e:
    print("ERROR during PowerPoint export:")
    traceback.print_exc()